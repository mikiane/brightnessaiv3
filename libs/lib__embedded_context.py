

# ----------------------------------------------------------------------------
# Project: Semantic Search Module for the Alter Brain project
# File:    lib__embedded_context.py
# 
# This lib is the Semantic Search Module for the Alter Brain project. It implements a 
# system for understanding and processing natural language to facilitate 
# information retrieval based on semantics rather than traditional keyword-based search.
# 
# Author:  Michel Levy Provencal
# Brightness.ai - 2023 - contact@brightness.fr
# ----------------------------------------------------------------------------

import pandas as pd
import os
import csv
from openai import OpenAI
from transformers import GPT2TokenizerFast
import random
import numpy as np
import sys
import time
import requests
import os.path
import PyPDF2
import docx
import json
import pptx
import xml.etree.ElementTree as ET
from bs4 import BeautifulSoup
import pytesseract
from PIL import Image
from openpyxl import load_workbook
from bs4 import BeautifulSoup
from markdownify import markdownify as md
from urllib.parse import urlparse, urljoin

# Import centralized configuration
from libs import lib__config as config

os.environ["TOKENIZERS_PARALLELISM"] = "false"

# Modèle par défaut via config
DEFAULT_MODEL = config.DEFAULT_MODEL
model = DEFAULT_MODEL

# -------------------- TOOLS --------------------

def generate_unique_filename(prefix, suffix):
    random_number = random.randint(1, 9999)
    return f"{prefix}_{random_number}.{suffix}"

# PDF to text
def convert_pdf_to_text(file_path):
    with open(file_path, "rb") as file:
        pdf_reader = PyPDF2.PdfReader(file)
        text = ""
        for page_num in range(len(pdf_reader.pages)):
            page = pdf_reader.pages[page_num]
            text += page.extract_text()
        return text

# DOCX to text
def convert_docx_to_text(file_path):
    doc = docx.Document(file_path)
    text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
    return text

# CSV to text
def convert_csv_to_text(file_path):
    with open(file_path, "r") as file:
        csv_reader = csv.reader(file)
        text = "\n".join([",".join(row) for row in csv_reader])
        return text

# JSON to text
def convert_json_to_text(file_path):
    with open(file_path, "r") as file:
        data = json.load(file)
        text = json.dumps(data, indent=4)
        return text

# Excel to text
def convert_excel_to_text(file_path):
    workbook = load_workbook(file_path)
    text = ""
    for sheet in workbook:
        for row in sheet.values:
            text += ",".join([str(cell) for cell in row])
            text += "\n"
    return text

# PPTX to text
def convert_pptx_to_text(file_path):
    presentation = pptx.Presentation(file_path)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text += paragraph.text
                    text += "\n"
    return text

# XML to text
def convert_xml_to_text(file_path):
    tree = ET.parse(file_path)
    root = tree.getroot()
    text = ET.tostring(root, encoding="utf-8", method="text").decode("utf-8")
    return text

# HTML to text
def convert_html_to_text(file_path):
    with open(file_path, "r") as file:
        soup = BeautifulSoup(file, "html.parser")
        text = soup.get_text()
        return text

# Image to text
def convert_image_to_text(file_path):
    image = Image.open(file_path)
    text = pytesseract.image_to_string(image, lang="eng")
    return text

# Text file passthrough
def convert_text_to_text(file_path):
    with open(file_path, "r", encoding="utf-8") as file:
        text = file.read()
    return text

# Concat folder files to text
def concat_files_in_text(path):
    files = []
    for f in os.listdir(path):
        full_path = os.path.join(path, f)
        if os.path.isfile(full_path):
            files.append(full_path)
    texts = []
    for file in files:
        with open(file, 'r') as f:
            file_content = f.read()
            texts.append(file_content)
    return ' '.join(texts)

# Split text into blocks
def split_text_into_blocks(text, limit=4000):
    blocks = []
    current_block = ""
    words = text.split()
    for word in words:
        if len(current_block + word) + 1 < limit:
            current_block += word + " "
        else:
            last_delimiter_index = max(current_block.rfind(". "), current_block.rfind("\n"))
            if last_delimiter_index == -1:
                blocks.append(current_block.strip())
                current_block = word + " "
            else:
                delimiter = current_block[last_delimiter_index]
                blocks.append(current_block[:last_delimiter_index + (1 if delimiter == '.' else 0)].strip())
                current_block = current_block[last_delimiter_index + (2 if delimiter == '.' else 1):].strip() + " " + word + " "
    if current_block.strip():
        blocks.append(current_block.strip())
    return blocks

# Write blocks to csv
def write_blocks_to_csv(blocks, path, filename):
    with open(path + filename, "w", newline="", encoding="utf-8") as csvfile:
        csvwriter = csv.writer(csvfile, delimiter=',', quotechar='"', quoting=csv.QUOTE_MINIMAL)
        csvwriter.writerow(['Datas'])
        for block in blocks:
            csvwriter.writerow([block])

# OpenAI embeddings

def get_embedding(text, engine="text-embedding-ada-002"):
    text = text.replace("\n", " ")
    client = OpenAI(api_key=config.OPENAI_API_KEY)
    response = client.embeddings.create(input=[text], model=engine)
    response_dict = response.model_dump()
    return response_dict['data'][0]['embedding']

# Create embeddings from CSV
def create_embeddings(path, filename):
    tokenizer = GPT2TokenizerFast.from_pretrained("gpt2")
    df_full = pd.read_csv(path + filename, sep=';', on_bad_lines='skip', encoding='utf-8')
    df = df_full[['Datas']]
    df = df.dropna()
    df['n_tokens'] = df.Datas.apply(lambda x: len(tokenizer.encode(x)))
    df = df[df.n_tokens < 8000].tail(2000)
    df['ada_embedding'] = df.Datas.apply(get_embedding)
    df.to_csv(path + "emb_" + filename, index=False)
    return path

# Read and process CSV for search
def read_and_process_csv(index_filename):
    df = pd.read_csv(index_filename)
    df['ada_embedding'] = df.ada_embedding.apply(eval).apply(np.array)
    return df

# Search helpers
def get_search_vector(text):
    return get_embedding(text)

def find_similar_rows(df, searchvector, n_results):
    df['similarities'] = df.ada_embedding.apply(lambda x: np.dot(x, searchvector))
    res = df.sort_values('similarities', ascending=False).head(n_results)
    return res

def validate_and_get_combined(res):
    if 'Datas' not in res.columns:
        raise ValueError("La colonne 'Datas' n'existe pas dans le DataFrame")
    if res.empty:
        raise ValueError("Le DataFrame est vide")
    if res.index.dtype != 'int64':
        raise ValueError("L'index du DataFrame n'est pas de type entier")
    return '\n'.join(res['Datas'].values)

# Indexing functions
def create_text_folder(folder_path):
    source_folder = folder_path
    destination_folder = folder_path + "text_tmp"
    supported_formats = {
        ".pdf": convert_pdf_to_text,
        ".docx": convert_docx_to_text,
        ".csv": convert_csv_to_text,
        ".json": convert_json_to_text,
        ".xls": convert_excel_to_text,
        ".xlsx": convert_excel_to_text,
        ".pptx": convert_pptx_to_text,
        ".xml": convert_xml_to_text,
        ".html": convert_html_to_text,
        ".jpg": convert_image_to_text,
        ".jpeg": convert_image_to_text,
        ".png": convert_image_to_text,
        ".txt": convert_text_to_text
    }
    if not os.path.exists(destination_folder):
        os.makedirs(destination_folder)
    for file_name in os.listdir(source_folder):
        source_file_path = os.path.join(source_folder, file_name)
        if os.path.isfile(source_file_path):
            file_name_without_ext, file_ext = os.path.splitext(file_name)
            if file_ext.lower() in supported_formats:
                converter = supported_formats[file_ext.lower()]
                text = converter(source_file_path)
                destination_file_name = generate_unique_filename(file_name_without_ext, "txt")
                destination_file_path = os.path.join(destination_folder, destination_file_name)
                with open(destination_file_path, "w", encoding="utf-8") as file:
                    file.write(text)
    return str(destination_folder)

# Build index from folder
def build_index(folder_path):
    text_folder_path = create_text_folder(folder_path)
    text = concat_files_in_text(text_folder_path)
    blocks = split_text_into_blocks(text, limit=4000)
    write_blocks_to_csv(blocks, folder_path, 'index.csv')
    brain_id = create_embeddings(folder_path, 'index.csv')
    return brain_id

# Build index from URL
def build_index_url(url):
    index_text = get_text_from_url(url)
    timestamp = time.strftime("%Y%m%d-%H%M%S")
    folder_path = "datas/" + timestamp + "/"
    os.makedirs(folder_path, exist_ok=True)
    with open(folder_path + 'url_index.txt', 'w', encoding='utf-8') as f:
        f.write(index_text)
    build_index(folder_path)
    return timestamp

# Retrieve context
def find_context(text, index_filename, n_results=5):
    if not os.path.exists(index_filename):
        return ""
    df = read_and_process_csv(index_filename)
    searchvector = get_search_vector(text)
    res = find_similar_rows(df, searchvector, n_results)
    return validate_and_get_combined(res)

# Query OpenAI with context
def query_extended_llm(text, index_filename, model=DEFAULT_MODEL):
    context = find_context(text, index_filename)
    client = OpenAI(api_key=config.OPENAI_API_KEY)
    attempts = 0
    prompt = "Context : " + context + "\n\n" + "Query : " + text
    while attempts < 10:
        try:
            response = client.chat.completions.create(
                model=model,
                messages=[{"role": "user", "content": prompt}],
            )
            message = response.choices[0].message.content
            return message.strip()
        except Exception as e:
            attempts += 1
            print(f"Erreur : {type(e).__name__} - {e}. Nouvel essai dans 5 secondes...")
            time.sleep(int(attempts) * 2)
    print("Erreur : Echec de la création de la completion après 10 essais")
    raise RuntimeError("LLM failure")
